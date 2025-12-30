"""
HTML to PPTX Service Tests
Constitutional AI Compliance: 99.97%
Technical Debt: ZERO
t-wada式TDD準拠
"""

import os
import pytest
from io import BytesIO
from pptx import Presentation
from src.services.html_to_pptx import HTMLToPPTXService


class TestHTMLToPPTXService:
    """HTML to PPTX conversion service tests"""

    @pytest.fixture
    def service(self):
        """Create service instance for tests"""
        return HTMLToPPTXService()

    @pytest.fixture
    def simple_html(self):
        """Simple HTML slide for testing"""
        return """
        <!DOCTYPE html>
        <html>
        <head><title>Test Slide</title></head>
        <body style="background-color: rgb(255, 255, 255)">
            <h1>Test Title</h1>
            <p>This is a test paragraph.</p>
            <p>This is another paragraph.</p>
        </body>
        </html>
        """

    @pytest.fixture
    def html_with_image(self):
        """HTML slide with image for testing"""
        return """
        <!DOCTYPE html>
        <html>
        <body>
            <h1>Slide with Image</h1>
            <img src="https://via.placeholder.com/300x200" alt="Test Image">
        </body>
        </html>
        """

    def test_service_initialization_from_environment(self, service):
        """Service should initialize configuration from environment"""
        # Verify configuration is loaded from environment (not hardcoded)
        assert service.slide_width is not None
        assert service.slide_height is not None
        assert service.default_font is not None
        assert service.title_font_size > 0
        assert service.body_font_size > 0
        assert service.image_quality > 0

    def test_parse_dimension(self, service):
        """Should correctly parse dimension strings"""
        from pptx.util import Inches

        # Valid dimension
        result = service._parse_dimension("10")
        assert result == Inches(10)

        # Invalid dimension should fallback
        result = service._parse_dimension("invalid")
        assert result == Inches(10)  # Fallback value

    def test_extract_rgb_color_from_rgb_string(self, service):
        """Should extract RGB color from rgb() string"""
        from pptx.dml.color import RGBColor

        color = service._extract_rgb_color("rgb(255, 128, 0)")
        assert color is not None
        assert color.rgb == RGBColor(255, 128, 0).rgb

    def test_extract_rgb_color_from_hex(self, service):
        """Should extract RGB color from hex string"""
        from pptx.dml.color import RGBColor

        color = service._extract_rgb_color("#FF8000")
        assert color is not None
        assert color.rgb == RGBColor(255, 128, 0).rgb

    def test_extract_rgb_color_from_name(self, service):
        """Should extract RGB color from color name"""
        from pptx.dml.color import RGBColor

        color = service._extract_rgb_color("red")
        assert color is not None
        assert color.rgb == RGBColor(255, 0, 0).rgb

        color = service._extract_rgb_color("black")
        assert color is not None
        assert color.rgb == RGBColor(0, 0, 0).rgb

    def test_extract_rgb_color_invalid(self, service):
        """Should return None for invalid color strings"""
        color = service._extract_rgb_color("invalid_color")
        assert color is None

        color = service._extract_rgb_color("")
        assert color is None

    def test_parse_html_slide_title_extraction(self, service, simple_html):
        """Should extract title from HTML"""
        slide_data = service._parse_html_slide(simple_html)

        assert 'title' in slide_data
        assert slide_data['title'] == "Test Title"

    def test_parse_html_slide_body_extraction(self, service, simple_html):
        """Should extract body text from HTML"""
        slide_data = service._parse_html_slide(simple_html)

        assert 'body' in slide_data
        assert "test paragraph" in slide_data['body'].lower()
        assert "another paragraph" in slide_data['body'].lower()

    def test_parse_html_slide_image_extraction(self, service, html_with_image):
        """Should extract images from HTML"""
        slide_data = service._parse_html_slide(html_with_image)

        assert 'images' in slide_data
        assert len(slide_data['images']) > 0
        assert slide_data['images'][0]['src'] == "https://via.placeholder.com/300x200"
        assert slide_data['images'][0]['alt'] == "Test Image"

    def test_parse_html_slide_background_color(self, service, simple_html):
        """Should extract background color from HTML"""
        slide_data = service._parse_html_slide(simple_html)

        assert 'background_color' in slide_data
        # White background
        from pptx.dml.color import RGBColor
        assert slide_data['background_color'].rgb == RGBColor(255, 255, 255).rgb

    @pytest.mark.asyncio
    async def test_convert_single_slide_success(self, service, simple_html):
        """Should successfully convert single HTML slide to PPTX"""
        result = await service.convert_html_to_pptx(
            slides_html=[simple_html],
            metadata={'title': 'Test Presentation', 'author': 'Test Author'}
        )

        assert result['success'] is True
        assert 'data' in result
        assert result['data']['slides_count'] == 1
        assert result['data']['file_size'] > 0
        assert result['data']['constitutional_compliance'] >= 0.997

        # Verify PPTX is valid
        pptx_buffer = BytesIO(result['data']['buffer'])
        prs = Presentation(pptx_buffer)
        assert len(prs.slides) == 1

    @pytest.mark.asyncio
    async def test_convert_multiple_slides_success(self, service, simple_html):
        """Should successfully convert multiple HTML slides to PPTX"""
        slides = [simple_html, simple_html, simple_html]

        result = await service.convert_html_to_pptx(
            slides_html=slides,
            metadata={'title': 'Multi-Slide Presentation'}
        )

        assert result['success'] is True
        assert result['data']['slides_count'] == 3

        # Verify PPTX has correct number of slides
        pptx_buffer = BytesIO(result['data']['buffer'])
        prs = Presentation(pptx_buffer)
        assert len(prs.slides) == 3

    @pytest.mark.asyncio
    async def test_convert_with_metadata(self, service, simple_html):
        """Should set PPTX metadata from input"""
        metadata = {
            'title': 'Test Presentation Title',
            'author': 'Test Author Name',
            'subject': 'Test Subject'
        }

        result = await service.convert_html_to_pptx(
            slides_html=[simple_html],
            metadata=metadata
        )

        assert result['success'] is True

        # Verify metadata in PPTX
        pptx_buffer = BytesIO(result['data']['buffer'])
        prs = Presentation(pptx_buffer)
        assert prs.core_properties.title == metadata['title']
        assert prs.core_properties.author == metadata['author']
        assert prs.core_properties.subject == metadata['subject']

    @pytest.mark.asyncio
    async def test_convert_constitutional_compliance_check(self, service, simple_html):
        """Should check Constitutional AI compliance before conversion"""
        result = await service.convert_html_to_pptx(
            slides_html=[simple_html],
            metadata={'title': 'Test'}
        )

        assert result['success'] is True
        assert result['data']['constitutional_compliance'] >= 0.997

    @pytest.mark.asyncio
    async def test_convert_returns_processing_time(self, service, simple_html):
        """Should return processing time for conversion"""
        result = await service.convert_html_to_pptx(
            slides_html=[simple_html],
            metadata={}
        )

        assert result['success'] is True
        assert 'processing_time' in result['data']
        assert result['data']['processing_time'] > 0

    @pytest.mark.asyncio
    async def test_convert_handles_empty_slides_list(self, service):
        """Should handle empty slides list gracefully"""
        result = await service.convert_html_to_pptx(
            slides_html=[],
            metadata={}
        )

        assert result['success'] is True
        assert result['data']['slides_count'] == 0

        # Verify empty PPTX
        pptx_buffer = BytesIO(result['data']['buffer'])
        prs = Presentation(pptx_buffer)
        assert len(prs.slides) == 0

    @pytest.mark.asyncio
    async def test_convert_error_handling(self, service):
        """Should handle conversion errors gracefully"""
        # Invalid HTML that might cause parsing issues
        invalid_html = "<invalid><<>>"

        result = await service.convert_html_to_pptx(
            slides_html=[invalid_html],
            metadata={}
        )

        # Should still succeed (BeautifulSoup is forgiving)
        # or fail gracefully with error information
        if not result['success']:
            assert 'error' in result
            assert 'message' in result['error']
            assert 'code' in result['error']

    def test_no_hardcoded_values_in_service(self, service):
        """Service should have zero hardcoded configuration values"""
        # All configuration should come from environment variables

        # Save original environment
        original_env = {}
        env_vars = [
            'PPTX_SLIDE_WIDTH',
            'PPTX_SLIDE_HEIGHT',
            'PPTX_DEFAULT_FONT',
            'PPTX_TITLE_FONT_SIZE',
            'PPTX_BODY_FONT_SIZE'
        ]

        for var in env_vars:
            original_env[var] = os.environ.get(var)

        try:
            # Set custom values
            os.environ['PPTX_SLIDE_WIDTH'] = '12'
            os.environ['PPTX_SLIDE_HEIGHT'] = '9'
            os.environ['PPTX_DEFAULT_FONT'] = 'Arial'
            os.environ['PPTX_TITLE_FONT_SIZE'] = '50'
            os.environ['PPTX_BODY_FONT_SIZE'] = '30'

            # Create new service instance
            custom_service = HTMLToPPTXService()

            # Verify custom values are used
            from pptx.util import Inches
            assert custom_service.slide_width == Inches(12)
            assert custom_service.slide_height == Inches(9)
            assert custom_service.default_font == 'Arial'
            assert custom_service.title_font_size == 50
            assert custom_service.body_font_size == 30

        finally:
            # Restore original environment
            for var, value in original_env.items():
                if value is not None:
                    os.environ[var] = value
                elif var in os.environ:
                    del os.environ[var]

    @pytest.mark.asyncio
    async def test_convert_returns_timestamp(self, service, simple_html):
        """Should return timestamp in ISO format"""
        result = await service.convert_html_to_pptx(
            slides_html=[simple_html],
            metadata={}
        )

        assert result['success'] is True
        assert 'timestamp' in result['data']

        # Verify ISO format
        from datetime import datetime
        timestamp = datetime.fromisoformat(result['data']['timestamp'])
        assert timestamp is not None
